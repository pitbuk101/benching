"""Export output of negotiation factory to a PowerPoint slide."""

import base64
import re
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ada.use_cases.negotiation_factory.negotiation_ppt_gen_utils import (
    add_circle_on_line,
    check_if_grouped_shape_contains_text,
    data_preprocessing_for_ppt,
    delete_extra_slides,
    delete_unwanted_nego_objective_sliders,
    find_key_recursive,
    find_keys_and_replace_text_in_shape,
    replace_text_in_paragraph,
    validate_nego_data,
)
from ada.utils.config.config_loader import read_config
from ada.utils.logs.logger import get_logger
from ada.utils.logs.time_logger import log_time

ppt_generation_config = read_config("use-cases.yml")["nf_ppt_generation"]
log = get_logger("nf-ppt-generation")

TEMPLATE = Path(Path(__file__).parents[4], "data/ppt_generation/templates/template_nego.pptx")
MAX_OPPORTUNITIES_PER_SLIDE = ppt_generation_config["MAX_OPPORTUNITIES_PER_SLIDE"]
MAX_NEGO_OBJECTIVES_PER_SLIDE = ppt_generation_config["MAX_NEGO_OBJECTIVES_PER_SLIDE"]
CHART_TITLE_ANNUAL_SPEND = ppt_generation_config["CHART_TITLE_ANNUAL_SPEND"]
CHART_TITLE_SKU_SPEND = ppt_generation_config["CHART_TITLE_SKU_SPEND"]
TITLE_SLIDE_INDEX = ppt_generation_config["TITLE_SLIDE_INDEX"]
NEGO_SUMMARY_SLIDE_INDICES = ppt_generation_config["NEGO_SUMMARY_SLIDE_INDICES"]
OPPORTUNITIES_SLIDE_INDICES = ppt_generation_config["OPPORTUNITIES_SLIDE_INDICES"]
CARROT_STICKS_SLIDE_INDEX = ppt_generation_config["CARROT_STICKS_SLIDE_INDEX"]
FINAL_ROUND_OUTCOME_SLIDE_INDICES = ppt_generation_config["FINAL_ROUND_OUTCOME_SLIDE_INDICES"]
ADDITIONAL_SLIDES_TO_DELETE = ppt_generation_config["ADDITIONAL_SLIDES_TO_DELETE"]

UNIT_MAP = {"PERCENTAGE": "%", "NUMBER": "", "DAYS": "days", "MONTHS": "months", "TEXT": "%"}


def update_title_slide(title_slide, negotiation_data):
    """
    Populate the title slide with the negotiation data.
    """
    for shape in title_slide.shapes:
        if shape.has_text_frame:
            find_keys_and_replace_text_in_shape(shape, negotiation_data)


def populate_objective_slider_group(
    grouped_shape: GroupShape,
    nego_summary_slide: Slide,
    negotiation_data: dict[str, Any],
):
    """
    Populate the objective slider group with the negotiation data.
    The text replacements are done and dot to the line is added
    args:
        grouped_shape (GroupShape): The grouped shape to populate.
        nego_summary_slide (Slide): The slide containing the grouped shape.
        negotiation_data (dict[str, Any]): The negotiation data to populate the shape with.
    """
    objective_id = None
    for shape in grouped_shape.shapes:
        if shape.has_text_frame:
            match = re.search(r"<objective(\d+)objectiveType>", shape.text_frame.text)
            if match:
                objective_id = int(match.group(1))
            find_keys_and_replace_text_in_shape(shape, negotiation_data)

    if objective_id is not None:
        for shape in grouped_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.LINE and shape.line.dash_style is None:
                unit = find_key_recursive(negotiation_data, f"objective{objective_id}unit")
                unit_str = UNIT_MAP.get(unit, "")
                min_value = find_key_recursive(
                    negotiation_data,
                    f"objective{objective_id}currentValue",
                )
                try:
                    min_value = float(min_value)
                except ValueError:
                    min_value = None

                max_value = find_key_recursive(negotiation_data, f"objective{objective_id}target")
                try:
                    max_value = float(max_value)
                except ValueError:
                    max_value = None

                current_val = find_key_recursive(
                    negotiation_data,
                    f"objective{objective_id}currentOffer",
                )
                if current_val:
                    try:
                        current_val = float(current_val)
                    except ValueError:
                        current_val = None
                else:
                    # user has not entered the current offer so assume it is at lowest value
                    current_val = min_value

                if min_value is not None and max_value is not None and current_val is not None:
                    add_circle_on_line(
                        nego_summary_slide,
                        shape,
                        min_value,
                        max_value,
                        current_val,
                        unit_str=unit_str,
                    )


@log_time
def update_negotiation_summary_slides(
    nego_summary_slides: list[Slide],
    negotiation_data: dict[str, Any],
) -> None:
    """
    Populate the negotiation summary slides with the negotiation data.
    args:
        nego_summary_slides (list[Slide]): The slides to populate.
        negotiation_data (dict[str, Any]): The negotiation data to populate the slides with.
    return: None
    """
    if len(negotiation_data["objectives"]) <= 3:
        last_slide_id = 0
    else:
        last_slide_id = 1
    for i, nego_summary_slide in enumerate(nego_summary_slides):
        if i == last_slide_id:
            delete_unwanted_nego_objective_sliders(
                nego_summary_slide,
                len(negotiation_data["objectives"]),
            )

        for shape in nego_summary_slide.shapes:
            if shape.has_chart and shape.chart.has_title:
                title = shape.chart.chart_title.text_frame.text.strip()
                chart = shape.chart
                if title == CHART_TITLE_ANNUAL_SPEND:
                    chart_data = CategoryChartData()
                    chart_data.categories = negotiation_data["annual_spend_chart_data"][
                        "categories"
                    ]
                    chart_data.add_series(
                        negotiation_data["annual_spend_chart_data"]["series"][0]["name"],
                        tuple(negotiation_data["annual_spend_chart_data"]["series"][0]["values"]),
                    )
                    chart.replace_data(chart_data)
                    new_title = (
                        title
                        + " ("
                        + negotiation_data["annual_spend_chart_data"]["annualSpendcurrencySymbol"]
                        + ")"
                    )
                    replace_text_in_paragraph(
                        shape.chart.chart_title.text_frame.paragraphs[0],
                        re.compile(CHART_TITLE_ANNUAL_SPEND),
                        new_title,
                        match_multiple=False,
                    )

                elif title == CHART_TITLE_SKU_SPEND:
                    chart_data = CategoryChartData()
                    chart_data.categories = negotiation_data["top_sku_chart_data"]["categories"]
                    chart_data.add_series(
                        negotiation_data["top_sku_chart_data"]["series"][0]["name"],
                        tuple(negotiation_data["top_sku_chart_data"]["series"][0]["values"]),
                    )
                    chart.replace_data(chart_data)
                    new_title = (
                        title
                        + " ("
                        + negotiation_data["top_sku_chart_data"]["topSkuCurrencySymbol"]
                        + ")"
                    )
                    replace_text_in_paragraph(
                        shape.chart.chart_title.text_frame.paragraphs[0],
                        re.compile(CHART_TITLE_SKU_SPEND),
                        new_title,
                        match_multiple=False,
                    )

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and check_if_grouped_shape_contains_text(
                shape,
                "objectiveType>",
            ):
                populate_objective_slider_group(shape, nego_summary_slide, negotiation_data)

            if shape.has_text_frame:
                find_keys_and_replace_text_in_shape(shape, negotiation_data)


@log_time
def update_opportunities_slide(opportunities_slides: list[Slide], negotiation_data: dict[str, Any]):
    """
    Populate the opportunities slide with the negotiation data.
    args:
        opportunities_slide (Slide): The slide to populate.
        negotiation_data (dict[str, Any]): The negotiation data to populate the slide with.
    returns:
        None
    """
    column_names = ["Analytics Name", "Insight", "Opportunity Value"]
    table_data = [
        [
            i["analyticsName"],
            i["insight"],
            i.get("opportunity_value_str", ""),
            i.get("opportunity_value", None),
        ]
        for i in negotiation_data["insights"]
    ]
    table_data = sorted(
        table_data,
        key=lambda x: (x[3] is None, -x[3] if x[3] is not None else float("inf")),
    )
    table_data = [[i[0], i[1], i[2]] for i in table_data]

    table_data_chunked = []
    for i in range(0, len(table_data), MAX_OPPORTUNITIES_PER_SLIDE):
        end_index = i + MAX_OPPORTUNITIES_PER_SLIDE
        chunk_data = [column_names] + table_data[i:end_index]
        table_data_chunked.append(chunk_data)

    for i, opportunities_slide in enumerate(opportunities_slides):
        x_coordinate, y_coordinate, c_x, c_y = Inches(1), Inches(2), Inches(4), Inches(1.5)
        table_shape = opportunities_slide.shapes.add_table(
            len(table_data_chunked[i]),
            len(table_data_chunked[i][0]),
            x_coordinate,
            y_coordinate,
            c_x,
            c_y,
        )

        for j, row in enumerate(table_data_chunked[i]):
            for k, cell in enumerate(row):
                if cell is None:
                    cell = ""
                table_shape.table.cell(j, k).text = str(cell)
                for paragraph in table_shape.table.cell(j, k).text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
        # set column widths
        table_shape.table.columns[0].width = Inches(2.0)
        table_shape.table.columns[1].width = Inches(8.0)


def update_carrots_sticks_slide(
    carrots_sticks_slide: Slide, negotiation_data: dict[str, Any]
) -> None:
    """
    Populate carrots and sticks slide with the negotiation data.
    args:
        carrots_sticks_slide (Slide): The slide to populate.
        negotiation_data (dict[str, Any]): The negotiation data to populate the slide with.
    returns:
        None
    """
    for shape in carrots_sticks_slide.shapes:
        if shape.has_text_frame:
            find_keys_and_replace_text_in_shape(shape, negotiation_data)


@log_time
def generate_nego_ppt(negotiation_data: dict[str, Any]) -> bytes:
    """
    Generates a PowerPoint slide containing the negotiation data.

    Args:
        negotiation_data (dict[str,Any]): The negotiation data to be displayed on the slide.

    Returns:
        bytes: A byte array representing the PowerPoint file containing the generated slide.

    """

    data_preprocessing_for_ppt(negotiation_data)
    template = Presentation(TEMPLATE)

    updated_slides_ids = delete_extra_slides(
        template,
        len(negotiation_data["objectives"]),
        len(negotiation_data["insights"]),
    )
    title_slide = template.slides[0]

    negotiation_summary_template_slides = []
    for i in NEGO_SUMMARY_SLIDE_INDICES:
        revised_slide_index = updated_slides_ids[i]
        if revised_slide_index != -1:
            negotiation_summary_template_slides.append(template.slides[revised_slide_index])

    opportunities_slides = []
    for i in OPPORTUNITIES_SLIDE_INDICES:
        revised_slide_index = updated_slides_ids[i]
        if revised_slide_index != -1:
            opportunities_slides.append(template.slides[revised_slide_index])

    # Title Slide
    update_title_slide(title_slide, negotiation_data)

    # Nego Summary slides
    update_negotiation_summary_slides(negotiation_summary_template_slides, negotiation_data)

    # Opportunities Slide
    update_opportunities_slide(opportunities_slides, negotiation_data)

    # Carrots & Sticks Slide
    carrots_sticks_slide_id = updated_slides_ids[ppt_generation_config["CARROT_STICKS_SLIDE_INDEX"]]
    if carrots_sticks_slide_id != -1:
        carrots_sticks_slide = template.slides[carrots_sticks_slide_id]
        update_carrots_sticks_slide(carrots_sticks_slide, negotiation_data)

    # Save the presentation to a byte buffer
    buffer = BytesIO()
    template.save(buffer)

    buffer.seek(0)  # Rewind the buffer
    file_bytes = buffer.getvalue()  # write bytes array in dev

    # write bytes array in dev
    return file_bytes


@log_time
def run_nf_ppt_generation(
    json_data: dict[str, Any],
) -> dict[str, Any]:
    """
    Run the complete process of generating a
    NF PowerPoint presentation based on the provided JSON data.

    Args:
    json_file (str): Input JSON data as a string.

    Returns:
    dict[str, Any]: The generated PowerPoint file along with other metadata or an error response.
    """
    if not validate_nego_data(json_data):
        log.info("Invalid input data for NF PPT generation")
        response_data = {
            "file_name": "failed_ppt.pptx",
            "file_type": "pptx",
            "encoding_type": "base64",
            "char_encoding": "utf8",
            "is_encoded": True,
            "message": "PPT not generated due to insufficient data.",
            "response_type": "data_issue",
        }
    else:
        nego_ppt = generate_nego_ppt(json_data)

        if not nego_ppt:
            log.error("NF PPT file was not generated due to internal error.")
            raise ValueError("NF PPT file was not generated.")
        log.info(f"Returning PPT file, length: {len(nego_ppt)}")

        # Return the PowerPoint file as bytes
        response_data = {
            "file_name": "negotiations.pptx",
            "file_type": "pptx",
            "file_data": base64.b64encode(nego_ppt).decode("utf8"),
            "encoding_type": "base64",
            "char_encoding": "utf8",
            "is_encoded": True,
            "message": "Successfully PPT generated",
            "response_type": "general",
        }
    return response_data
