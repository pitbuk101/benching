"""Utility functions for Export output of negotiation factory to a PowerPoint slide."""

import json
import math
import re
from typing import Any

import pptx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ada.components.db.pg_connector import PGConnector
from ada.components.llm_models.generic_calls import generate_chat_response_with_chain
from ada.use_cases.negotiation_factory.prompts import get_opportunity_extraction_prompt
from ada.utils.config.config_loader import read_config
from ada.utils.logs.logger import get_logger
from ada.utils.logs.time_logger import log_time

ppt_generation_config = read_config("use-cases.yml")["nf_ppt_generation"]
log = get_logger("nf-ppt-generation")

UNIT_MAP = {"PERCENTAGE": "%", "NUMBER": "", "DAYS": "days", "MONTHS": "months", "TEXT": "%"}


def add_prefix_to_list_of_dicts(data: list[dict[str, Any]], prefix: str):
    """
    Add a prefix_txt + index to the keys of each dictionary in a list of dictionaries.

    Args:
        data (list[dict[str, Any]): The list of dictionaries to modify.
        prefix (str): The prefix to add to the keys.

    Returns:
        None
    """
    for i, item in enumerate(data):
        keys_to_replace = list(item.keys())
        for key in keys_to_replace:
            item[f"{prefix}{i + 1}{key}"] = item.pop(key)


def format_large_numbers(numbers: list[float]) -> tuple[list[float], str]:
    """
    Formats a list of numbers to a common scale based on the smallest number in the list.
    Returns a list of rounded numbers and the chosen scale.
    args:
        numbers (list[float]): A list of numbers to format.
    Returns:
        tuple[list[float], str]: A tuple containing the formatted numbers and the chosen scale.
    """
    if not numbers:
        return [], ""  # Handle empty list

    smallest_number = min(numbers)

    # Determine the scale based on the smallest number
    if smallest_number >= 1_000_000_000:  # Billion
        scale_divisor = 1_000_000_000
        scale = "Bn"
    elif smallest_number >= 1_000_000:  # Million
        scale_divisor = 1_000_000
        scale = "Mn"
    elif smallest_number >= 1_000:  # Thousand
        scale_divisor = 1_000
        scale = "K"
    else:  # Below thousand
        scale_divisor = 1
        scale = ""

    # Convert all numbers to the chosen scale
    rounded_numbers = [round(num / scale_divisor, 1) for num in numbers]

    return rounded_numbers, scale


def extract_annual_spend_chart_data(negotiation_data: dict[str, Any]) -> dict[str, Any]:
    """
    Extracts the data required to create the annual spend chart in python-pptx package chart format
    args:
        negotiation_data (dict[str, Any]): The negotiation data to extract the chart data from.
    Returns:
        dict[str, Any]: A dictionary containing the chart data.
    """
    spend_numbers = [
        negotiation_data["supplier"]["spendLastYear"],
        negotiation_data["supplier"]["spendYtd"],
    ]
    spend_numbers, spend_scale = format_large_numbers(spend_numbers)
    years = [
        str(int(negotiation_data["supplier"]["period"]) - 1),
        str(negotiation_data["supplier"]["period"]),
    ]
    chart_data = {
        "categories": years,
        "series": [{"name": "annual_spend", "values": spend_numbers}],
        "annualSpendcurrencySymbol": negotiation_data["supplier"]["currencySymbol"]
        + " "
        + spend_scale,
    }
    return chart_data


def extract_top_sku_chart_data(negotiation_data: dict[str, Any]) -> dict[str, Any]:
    """
    Extracts the data required to create the top SKU chart in python-pptx package chart format
    args:
        negotiation_data (dict[str, Any]): The negotiation data to extract the chart data from.
    Returns:
        dict[str, Any]: A dictionary containing the chart data.
    """
    if negotiation_data.get("skus"):
        top_sku_data = negotiation_data["skus"]
        top_sku_data = sorted(top_sku_data, key=lambda x: x["spend"], reverse=True)
        if len(top_sku_data) <= ppt_generation_config["MAX_SKU_CNT"]:
            top_sku_names = [sku["name"] for sku in top_sku_data]
            top_sku_spend = [sku["spend"] for sku in top_sku_data]
        else:
            top_sku_names = [
                sku["name"] for sku in top_sku_data[: ppt_generation_config["MAX_SKU_CNT"]]
            ]
            top_sku_names.append("Others")
            top_sku_spend = [
                sku["spend"] for sku in top_sku_data[: ppt_generation_config["MAX_SKU_CNT"]]
            ]
            top_sku_spend.append(
                sum(sku["spend"] for sku in top_sku_data[ppt_generation_config["MAX_SKU_CNT"] :]),
            )
        top_sku_names.append("Total")
        top_sku_spend.append(sum(sku["spend"] for sku in top_sku_data))

        top_sku_spend, sku_scale = format_large_numbers(top_sku_spend)
        chart_data = {
            "categories": top_sku_names,
            "series": [{"name": "top_sku_spend", "values": top_sku_spend}],
            "topSkuCurrencySymbol": negotiation_data["supplier"]["currencySymbol"]
            + " "
            + sku_scale,
        }
    else:
        chart_data = {
            "categories": [],
            "series": [{"name": "top_sku_spend", "values": []}],
            "topSkuCurrencySymbol": "",
        }
    return chart_data


def preprocess_objectives(negotiation_data: dict[str, Any]) -> None:
    """
    preprocess objectives in negotiation data
    args:
        negotiation_data (dict[str, Any]): The negotiation data to preprocess.
    return: None
    """
    if negotiation_data.get("objectives") is None:
        negotiation_data["objectives"] = []
    negotiation_data["objectivesCnt"] = len(negotiation_data["objectives"])

    # add units strings to objectives
    for objective in negotiation_data["objectives"]:
        objective["unitStr"] = UNIT_MAP.get(objective["unit"], "")
        objective["objectiveParam"] = objective["objectiveParam"].lower()

    # add prefix to objectives, tactics and sourcing approach
    add_prefix_to_list_of_dicts(negotiation_data["objectives"], "objective")


def preprocess_supplier(negotiation_data: dict[str, Any]) -> None:
    """
    preprocess supplier in negotiation data
    args:
        negotiation_data (dict[str, Any]): The negotiation data to preprocess.
    returns: None
    """
    # nego spend rounding
    nego_spend_number, nego_spend_scale = format_large_numbers([negotiation_data["negoSpend"]])
    if negotiation_data["supplier"]["currencyPosition"] == "prefix":
        negotiation_data["negoSpendRounded"] = (
            negotiation_data["supplier"]["currencySymbol"]
            + str(nego_spend_number[0])
            + " "
            + nego_spend_scale
        )
    else:
        negotiation_data["negoSpendRounded"] = (
            str(nego_spend_number[0])
            + " "
            + negotiation_data["supplier"]["currencySymbol"]
            + nego_spend_scale
        )

    # spend percentage
    negotiation_data["supplier"]["percentageSpendAcrossCategoryYtd"] = round(
        negotiation_data["supplier"]["percentageSpendAcrossCategoryYtd"] * 100,
        2,
    )
    negotiation_data["supplier"]["percentageSpendAcrossCategoryLastYear"] = round(
        negotiation_data["supplier"]["percentageSpendAcrossCategoryLastYear"] * 100,
        2,
    )
    negotiation_data["supplier"]["percentageSpendWhichIsSingleSourced"] = round(
        negotiation_data["supplier"]["percentageSpendWhichIsSingleSourced"] * 100,
        2,
    )


def preprocess_tones_and_tactics(negotiation_data: dict[str, Any]) -> None:
    """
    preprocess tones and tactics in negotiation data
    args: negotiation_data (dict[str, Any]): The negotiation data to preprocess.
    return: None
    """
    if negotiation_data.get("toneAndTactics") and negotiation_data["toneAndTactics"].get("tactics"):
        negotiation_data["tactics_list"] = [
            i.get("title", "") + " : " + i.get("description", "")
            for i in negotiation_data["toneAndTactics"]["tactics"]
        ]
    else:
        negotiation_data["tactics_list"] = ""

    if negotiation_data.get("toneAndTactics"):
        negotiation_data["tone"] = (
            negotiation_data["toneAndTactics"].get("title")
            + " : "
            + negotiation_data["toneAndTactics"].get("description")
        )
    else:
        negotiation_data["tone"] = ""


def get_top_suppliers(negotiation_data: dict[str, Any], supplier_cnt: int = 3) -> None:
    """
    Get the top suppliers based on spend YTD.
    args
    """
    # read data from DB for competitors
    pg_db_conn = PGConnector(tenant_id=negotiation_data["tenant_id"])

    top_suppliers = pg_db_conn.select_records_with_filter(
        table_name="supplier_profile_with_insights_with_objectives_view_with_saving",
        filtered_columns=["supplier_name"],
        filter_condition="period = (SELECT MAX(period) FROM "
        "supplier_profile_with_insights_with_objectives_view_with_saving)",
        num_records=supplier_cnt,
        order_by=("spend_ytd", "desc"),
    )
    negotiation_data["top3_suppliers"] = ", ".join([supplier[0] for supplier in top_suppliers])


def preprocess_insights(negotiation_data):
    """
    preprocess insights in negotiation data
    args:
        negotiation_data (dict[str, Any]): The negotiation data to preprocess.
    returns: None
    """
    # extract opportunity values for insights
    opportunity_values = extract_opportunity_values_from_insights(negotiation_data["insights"])
    for opportunity in negotiation_data["insights"]:
        if (
            opportunity_values.get(opportunity["id"])
            and opportunity_values[opportunity["id"]].get("opportunity_value")
            and opportunity_values[opportunity["id"]].get("opportunity_currency")
        ):
            scale_map = {"K": 1000, "M": 1000000, "Mn": 1000000, "B": 1000000000, "Bn": 1000000000}
            currency = opportunity_values[opportunity["id"]].get("opportunity_currency", "")
            currency = "" if currency is None else currency

            scale = opportunity_values[opportunity["id"]].get("opportunity_scale", "")
            scale = "" if scale is None else scale

            opportunity["opportunity_value_str"] = (
                str(opportunity_values[opportunity["id"]].get("opportunity_value"))
                + " "
                + scale
                + " "
                + currency
            )

            opportunity_value = opportunity_values[opportunity["id"]].get("opportunity_value")
            try:
                opportunity_value = float(opportunity_value)
                opportunity["opportunity_value"] = opportunity_value * scale_map.get(scale, 1)
            except ValueError:
                opportunity["opportunity_value"] = None


def preprocess_sourcing_approach(negotiation_data: dict[str, Any]) -> None:
    """
    preprocess sourcing approach in negotiation data
    """
    if negotiation_data.get("sourcingApproach"):
        negotiation_data["sourcingApproach_list"] = [
            i["type"] + " : " + i["value"] for i in negotiation_data["sourcingApproach"]
        ]
    else:
        negotiation_data["sourcingApproach_list"] = ""


def preprocess_carrots_sticks(negotiation_data: dict[str, Any]) -> None:
    """
    preprocess carrots and sticks in negotiation data
    args:   negotiation_data (dict[str, Any]): The negotiation data to preprocess.
    return: None
    """
    if negotiation_data.get("carrots"):
        negotiation_data["carrots_list"] = [
            i.get("title", "") + " : " + i.get("description", "")
            for i in negotiation_data.get("carrots", [])
        ]
    else:
        negotiation_data["carrots_list"] = ""

    if negotiation_data.get("sticks"):
        negotiation_data["sticks_list"] = [
            i.get("title", "") + " : " + i.get("description", "")
            for i in negotiation_data.get("sticks", [])
        ]
    else:
        negotiation_data["sticks_list"] = ""


@log_time
def data_preprocessing_for_ppt(negotiation_data: dict[str, Any]) -> None:
    """
    Each text box in template ppt slide has a unique identifier.
    Modify input data in place to match this identifier
    1) add prefix to objectives, tactics and sourcing approach
    """
    preprocess_objectives(negotiation_data)

    preprocess_supplier(negotiation_data)

    log.info(
        f"Generating NF PPT for {negotiation_data['objectivesCnt']} objectives,"
        f" {len(negotiation_data['skus'])} SKUs, {len(negotiation_data['insights'])} insights,"
        f" {len(negotiation_data['supplier'])} suppliers",
    )

    preprocess_tones_and_tactics(negotiation_data)

    preprocess_sourcing_approach(negotiation_data)

    # extract chart data
    negotiation_data["annual_spend_chart_data"] = extract_annual_spend_chart_data(negotiation_data)
    negotiation_data["top_sku_chart_data"] = extract_top_sku_chart_data(negotiation_data)

    get_top_suppliers(negotiation_data)

    # add missing data which needs to be shown blank
    fill_blank_keys = ["buyerAttractiveness", "categoryPositioning", "supplierPositioning"]
    for key in fill_blank_keys:
        if negotiation_data.get(key) is None:
            negotiation_data[key] = ""
        else:
            negotiation_data[key] = negotiation_data[key].title()

    preprocess_insights(negotiation_data)

    preprocess_carrots_sticks(negotiation_data)


def find_key_recursive(data: dict[str, Any], target_key: str):
    """
    Recursively searches for a key in a nested dictionary.

    Args:
        data (dict): The dictionary to search.
        target_key (str): The key to find.

    Returns:
        The value associated with the target key, or None if the key is not found.
    """
    if target_key in data:
        return data[target_key]

    for _, value in data.items():
        if isinstance(value, dict):
            result = find_key_recursive(value, target_key)
            if result is not None:
                return result
        elif isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    result = find_key_recursive(item, target_key)
                    if result is not None:
                        return result
    return None


def replace_text_in_paragraph(
    paragraph,
    regex: re.Pattern,
    replace_str: str,
    match_multiple: bool = False,
):
    """Replacing all matches for `regex` with `replace_str`. The replacement is done in-place.
    args:
        paragraph (pptx.text.text.Paragraph): The paragraph to search in.
        regex (re.Pattern): The regex pattern to search for.
        replace_str(str): The string to replace the match with.
        match_multiple (bool): Whether to replace all matches in a paragraph or just the first one.
    """
    # --- a paragraph may contain more than one match, loop until all are replaced ---
    while True:
        text = paragraph.text
        match = regex.search(text)
        if not match:
            break

        # --- when there's a match, we need to modify run.text for each run that
        # --- contains any part of the match-string.
        runs = iter(paragraph.runs)
        start, end = match.start(), match.end()

        # --- Skip over any leading runs that do not contain the match ---
        run = None
        for run in runs:
            run_len = len(run.text)
            if start < run_len:
                break
            start, end = start - run_len, end - run_len

        if run is not None:
            # --- Match starts somewhere in the current run. Replace match-str prefix
            # --- occurring in this run with entire replacement str.
            run_text = run.text
            run_len = len(run_text)
            run.text = f"{run_text[:start]}{replace_str}{run_text[end:]}"
            end -= run_len  # --- note this is run-len before replacement ---

        # --- Remove any suffix of match word that occurs in following runs. Note that
        # --- such a suffix will always begin at the first character of the run. Also
        # --- note a suffix can span one or more entire following runs.
        for run in runs:  # --- next and remaining runs, uses same iterator ---
            if end <= 0:
                break
            run_text = run.text
            run_len = len(run_text)
            run.text = run_text[end:]
            end -= run_len
        if not match_multiple:
            break


def find_keys_and_replace_text_in_shape(
    shape: pptx.shapes.autoshape.Shape,
    negotiation_data: dict[str, Any],
):
    """
    Find variables in the shape and replace it by looking up that key in the dictionary.
    Args:
        shape (pptx.slide.SlideShapes): The shape to search for variables.
        negotiation_data (dict[str,Any]): The dictionary to search for the variable.
    Returns:
        None
    """
    shape_text = shape.text_frame.text
    variable_names = re.findall(r"<(.*?)>", shape_text)
    for variable in variable_names:
        value = find_key_recursive(negotiation_data, variable)
        if value is not None:
            if isinstance(value, (str, int, float)):
                value = str(value)
                for paragraph in shape.text_frame.paragraphs:
                    replace_text_in_paragraph(
                        paragraph,
                        re.compile(f"<{variable}>"),
                        value,
                        match_multiple=True,
                    )

            elif isinstance(value, list):
                # Add the list as bullet points
                text_frame = shape.text_frame
                text_frame.clear()
                for line in value:
                    line = str(line)
                    if line.strip():
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = line.strip()
                        paragraph.level = 1
                        paragraph.font.size = Pt(ppt_generation_config["BULLET_PARA_FONT_SIZE"])


def delete_unwanted_nego_objective_sliders(nego_summary_slide: Slide, nego_objective_cnt: int):
    """
    Delete the unwanted negotiation objective sliders from the slide.
    There are 3 sliders on the slide.
    E.g. If there are 2 negotiation objectives, delete the 3rd slider.
    args:
        nego_summary_slide (Slide): The slide to delete the sliders from.
        nego_objective_cnt (int): The number of negotiation objectives to display.
        return: None
    """

    def should_delete_shape(shape: GroupShape, nego_objective_cnt: int) -> bool:
        for j in shape.shapes:
            if hasattr(j, "has_text_frame") and j.has_text_frame:
                match = re.search(r"<objective(\d+)objectiveType>", j.text_frame.text)
                if match:
                    objective_id = int(match.group(1))
                    if objective_id > nego_objective_cnt:
                        return True
        return False

    for shape in nego_summary_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and should_delete_shape(
            shape,
            nego_objective_cnt,
        ):
            shape.element.getparent().remove(shape.element)


def add_circle_on_line(
    slide: Slide,
    line_shape: Shape,
    min_value: float,
    max_value: float,
    value: float,
    radius: float = 0.1,
    color: str = "40E0D0",
    unit_str: str = "",
):
    """
    Adds a colored circle on a line shape at a specific value between a range.

    Args:
        slide (Slide): The slide containing the line shape.
        line_shape (Shape): The line shape to place the circle on.
        min_value (float): The minimum value of the range.
        max_value (float): The maximum value of the range.
        value (float): The value to place the circle at.
        radius (float): Radius of the circle in inches. Default is 0.1 inches.
        color (str): Hex color code for the circle. Default is red ('FF0000').

    Returns:
        Shape: The added circle shape.
    """
    # Ensure the value is within the range. If outside then set it to the min or max value
    if value < min_value:
        value = min_value
    elif value > max_value:
        value = max_value

    # Calculate the relative position of the value within the range
    relative_position = (value - min_value) / (max_value - min_value)

    # Get the start and end points of the line
    x_1, y_1 = line_shape.left, line_shape.top
    x_2 = line_shape.left + line_shape.width
    y_2 = line_shape.top + line_shape.height

    # Calculate the position of the circle based on the relative position
    circle_x = x_1 + relative_position * (x_2 - x_1)
    circle_y = y_1 + relative_position * (y_2 - y_1)

    # Add the circle to the slide
    circle_diameter = Inches(radius * 2)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=circle_x - circle_diameter / 2,
        top=circle_y - circle_diameter / 2,
        width=circle_diameter,
        height=circle_diameter,
    )

    # Set the circle color
    fill = circle.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)

    # Add a text box above the circle
    text_box_x = circle_x
    text_box_y = circle_y - int(circle_diameter) - Pt(10)  # Adjust 10 points above the circle
    text_box_width = int(circle_diameter)  # Same width as circle
    text_box_height = Pt(20)  # Set a fixed height for the text box

    text_box = slide.shapes.add_textbox(
        text_box_x,
        text_box_y,
        text_box_width,
        text_box_height,
    )
    text_frame = text_box.text_frame
    text_frame.text = str(value) + unit_str

    # Format the text
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)


def check_if_grouped_shape_contains_text(shape: GroupShape, text: str) -> bool:
    """
    Check if a grouped shape contains a text.
    args:
        shape (GroupShape): The grouped shape to search in.
        text (str): The text to search for.
    return:
        bool
    """
    for j in shape.shapes:
        if hasattr(j, "has_text_frame") and j.has_text_frame:
            if text in j.text_frame.text:
                return True
    return False


def delete_slides(prs, slide_indices):
    """
    Deletes slides at the specified indices in the presentation.
    """
    for i in reversed(slide_indices):
        r_id = prs.slides._sldIdLst[i].rId  # pylint: disable=protected-access
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]  # pylint: disable=protected-access


def update_indices_with_deletion(
    original_indices: list[int],
    deleted_indices: list[int],
) -> list[int]:
    """
    Return a new list of indices after deletion, replacing deleted indices with -1.

    Args:
        original_indices (list[int]): List of original indices.
        deleted_indices (list[int]): List of indices to be deleted.

    Returns:
        list[int]: Updated list of indices with -1 for deleted indices.
    """
    # Sort the list of deleted indices for proper processing
    deleted_indices = sorted(deleted_indices)

    # Initialize the result list with -1 for deleted indices
    result = [-1 if i in deleted_indices else i for i in original_indices]

    # Adjust the remaining indices to reflect deletions
    shift = 0  # Tracks the number of deletions processed so far
    for i, value in enumerate(result):
        if value == -1:
            shift += 1
        else:
            result[i] -= shift

    return result


def delete_extra_slides(
    template: pptx.Presentation,
    nego_objective_cnt: int,
    opportunities_cnt: int,
) -> list[int]:
    """
    Delete the extra slides from the template.
     Returns list which returns the updated slide indices after deletion.
    if a slide is deleted then values is -1 for that position in the list.
    args:
        template (pptx.Presentation): The template to delete the slides from.
        nego_objective_cnt (int): The number of negotiation objectives.
        opportunities_cnt (int): The number of opportunities.
    returns:
        list[int]: The updated slide indices after deletion
    """
    original_indices = list(range(len(template.slides)))
    slides_to_delete = []
    nego_summary_slides_to_keep = math.ceil(
        nego_objective_cnt / ppt_generation_config["MAX_NEGO_OBJECTIVES_PER_SLIDE"],
    )
    slides_to_delete.extend(
        ppt_generation_config["NEGO_SUMMARY_SLIDE_INDICES"][nego_summary_slides_to_keep:],
    )

    opportunities_slides_to_keep = math.ceil(
        opportunities_cnt / ppt_generation_config["MAX_OPPORTUNITIES_PER_SLIDE"],
    )
    slides_to_delete.extend(
        ppt_generation_config["OPPORTUNITIES_SLIDE_INDICES"][opportunities_slides_to_keep:],
    )

    slides_to_delete.extend(ppt_generation_config["ADDITIONAL_SLIDES_TO_DELETE"])
    delete_slides(template, slides_to_delete)

    updated_indices = update_indices_with_deletion(original_indices, slides_to_delete)
    return updated_indices


def validate_nego_data(json_data: dict[str, Any]) -> bool:
    """
    Validate the negotiation data. check if specific keys are present in the JSON data.
    args:
        json_data (dict[str]): The negotiation data to validate.
    return: bool
    """
    keys_to_check = ["objectives", "supplier"]
    for key in keys_to_check:
        if json_data.get(key) is None:
            return False
    return True


def extract_opportunity_values_from_insights(
    insights: list[dict[str, Any]],
) -> dict[str, str | Any]:
    """
    Extracts the opportunity values from the insights data.
    args:
        insights (list[dict[str, Any]]): The insights data to extract the opportunity values from.
    return: list[str]
    """
    prompt = get_opportunity_extraction_prompt(insights)
    ai_response = generate_chat_response_with_chain(
        prompt=prompt,
        model=ppt_generation_config["LLM"],
    )
    log.info("response from LLM: %s", ai_response)
    try:
        ai_response = json.loads(ai_response)
    except json.decoder.JSONDecodeError as exc:
        log.info("LLM generated output has mistake in json format %s", exc)

    # convert into dict for easy lookup
    extracted_opportunity_dict = {i.get("id"): i for i in ai_response if i}
    return extracted_opportunity_dict
