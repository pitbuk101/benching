"""PPT generation use case."""

import ast
import base64
import copy
import os
import re
from io import BytesIO
from typing import Any

import pptx
from pptx import Presentation
from pptx.util import Cm, Inches, Pt

from ada.components.llm_models.generic_calls import generate_chat_response
from ada.use_cases.ppt_generation.few_shot_examples import few_shot_examples
from ada.use_cases.ppt_generation.prompts import (
    get_prompt_chart_type,
    ppt_generate_cleaned_data_title_insights_prompt,
)
from ada.use_cases.ppt_generation.utils import (
    choose_model_based_on_token_length,
    find_index_of_element,
    get_data,
    get_pptx_template_path,
    populate_charts,
    to_superscript,
)
from ada.utils.config.config_loader import read_config
from ada.utils.logs.logger import get_logger
from ada.utils.logs.time_logger import log_time
from ada.utils.metrics.context_manager import (
    UseCase,
    set_tenant_and_intent_flow_context,
)

ppt_generation_config = read_config("use-cases.yml")["ppt_generation"]


log = get_logger("ppt-generation-v1")

model_conf = read_config("models.yml")


def generate_chart_type(
    user_question: str,
    data_from_answer: list[Any],
) -> str | None:
    """
    Generates the type of chart based on the user's question and the provided data.

    Args:
        user_question (str): A string containing the user's question or request for a chart.
        data_from_answer(list[Any]): Data that will be used to determine the appropriate chart type.
        The format and type of this data may vary.

    Returns:
        (str| None): The type of chart based on the user's question and the provided data.
    """
    prompt = get_prompt_chart_type(
        data_from_answer=data_from_answer,
        user_question=user_question,
        few_shot_examples=few_shot_examples,
    )

    # choose the right model based on token length
    model_name = choose_model_based_on_token_length(prompt[0]["content"])
    log.info(f"model_name: {model_name}")

    response = (
        generate_chat_response(
            messages=prompt,
            model=model_name,
        )
        if model_name
        else None
    )

    return response


def perform_rounding(
    cleaned_data: dict[str, Any],
    rounding: str,
    chart_type: str,
) -> dict[str, Any]:
    """
    Perform rounding on the cleaned data based on the chart type and rounding key
    args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        rounding (str): The rounding key to be applied on the data.
    returns:
        dict[str, Any]: The cleaned data with rounding applied.
    """
    round_digits_after_decimal = ppt_generation_config["round_digits_after_decimal"]
    rounded_data = {
        "categories": cleaned_data.get("categories", []),
        "original_series": cleaned_data.get("series", []),
        "series": [],
    }
    if rounding is not None:
        divide_by = get_rounding_division_number(rounding)

        for ser in cleaned_data.get("series", []):
            rounded_vals = [round(i / divide_by, round_digits_after_decimal) for i in ser[1]]
            rounded_tuple = (ser[0], tuple(rounded_vals))
            rounded_data["series"].append(rounded_tuple)
        if chart_type == "pie":
            rounded_data["original_series"] = rounded_data["series"]
    else:
        rounded_data["series"] = cleaned_data.get("series", [])
    return rounded_data


def adopt_rounding_for_chart_type(
    chart_type: str,
    rounding: str,
    cleaned_data: dict[str, Any],
    chart_tile: str,
) -> tuple[str, str]:
    """
    Adopt rounding for pie chart.
    IF the values are percentages then add another value to rounding as "%"
    Args:
        chart_type (str): The type of chart to be generated based on the user's question.
        rounding (str): The rounding key to be applied on the data.
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
    Returns:
        (str): The revised rounding key to be applied on the data.
    """
    rounding_options = ["Bn", "Mn", "K", "None"]
    if chart_type == "pie":
        sum_vals = sum(cleaned_data["series"][0][1])
        if 0.9 < sum_vals <= 1.1:
            rounding = "%"
    else:
        # if most of the values are becoming zero after rounding then update rounding
        zero_val_cnt = 0
        divide_by = get_rounding_division_number(rounding)
        for val in cleaned_data["series"][0][1][:5]:
            if round(val / divide_by, 0) == 0:
                zero_val_cnt += 1
        if zero_val_cnt > 1 and rounding in ["Bn", "Mn", "K"]:
            new_rounding = rounding_options[rounding_options.index(rounding) + 1]
            chart_tile = re.sub(r"\b" + rounding + r"\b", new_rounding, chart_tile)
            rounding = new_rounding

    return rounding, chart_tile


def get_rounding_division_number(rounding: str) -> float:
    """
    get the number to divide by as per rounding string
    args:
        rounding (str): The rounding key to be applied on the data.
    returns:
        (float): The number to divide by
    """
    if rounding is not None:
        if rounding == "K":
            divide_by = 1000.0
        elif rounding == "Mn":
            divide_by = 1000000.0
        elif rounding == "Bn":
            divide_by = 1000000000.0
        elif rounding == "%":
            divide_by = 0.01
        else:
            divide_by = 1.0
    else:
        divide_by = 1.0
    return divide_by


def sort_filter_data(chart_type: str, cleaned_data: dict[str, Any]) -> None:
    """
    Sort the data based on the chart type and remove the "#" values from the cleaned data
    args:
        chart_type (str): The type of chart to be generated based on the user's question.
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
    returns:
        None
    """
    if chart_type in ("bar", "pie"):
        sorted_categories, sorted_vals = zip(
            *sorted(
                zip(cleaned_data.get("categories", []), list(cleaned_data["series"][0][1])),
                key=lambda x: x[1],
                reverse=True,
            ),
        )
        # remove # values from categories and sorted_vals
        filtered_categories = []
        filtered_vals = []
        for cat, val in zip(sorted_categories, sorted_vals):
            if cat != "#":
                filtered_categories.append(cat)
                filtered_vals.append(val)

        cleaned_data["categories"] = list(filtered_categories)
        cleaned_data["series"][0] = (cleaned_data["series"][0][0], tuple(filtered_vals))

    elif chart_type == "stacked":
        # sort by first value in each tuple
        cleaned_data["series"] = sorted(cleaned_data["series"], key=lambda x: x[1][0], reverse=True)
        # remove # values from series
        cleaned_data["series"] = [i for i in cleaned_data["series"] if i[0] != "#"]


def concatenate_to_others(cleaned_data: dict[str, Any], chart_type: str) -> list[Any]:
    """
    Concatenate all the values more than 7 to Others. Return list of values concatenated to Others.
    Args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        chart_type (str): The type of chart to be generated based on the user's question.
    Returns:
        (list[Any]): The list of values concatenated
    """
    max_num_categories_to_show = ppt_generation_config["max_num_categories_to_show"]
    others_data: list[Any] = []
    if chart_type in ["pie", "bar"]:
        values_cnt = len(cleaned_data.get("categories", []))

        if values_cnt > max_num_categories_to_show + 1:
            others_categories = cleaned_data.get("categories", [])[max_num_categories_to_show:]
            others_values = cleaned_data["original_series"][0][1][max_num_categories_to_show:]
            others_sum = sum(cleaned_data["series"][0][1][max_num_categories_to_show:])
            others_sum_original = sum(
                cleaned_data["original_series"][0][1][max_num_categories_to_show:],
            )
            cleaned_data["categories"] = cleaned_data.get("categories", [])[
                :max_num_categories_to_show
            ] + [
                "Others",
            ]
            cleaned_data["series"][0] = (
                cleaned_data["series"][0][0],
                cleaned_data["series"][0][1][:max_num_categories_to_show] + (others_sum,),
            )
            cleaned_data["original_series"][0] = (
                cleaned_data["original_series"][0][0],
                cleaned_data["original_series"][0][1][:max_num_categories_to_show]
                + (others_sum_original,),
            )
            others_data = list(zip(others_categories, others_values))

    elif chart_type == "stacked":
        subcategories_cnt = len(cleaned_data.get("series", []))

        if subcategories_cnt > max_num_categories_to_show + 1:
            others_subcategories = cleaned_data.get("series", [])[max_num_categories_to_show:]
            others_values = []
            for i in range(len(others_subcategories[0][1])):
                others_values.append(sum(subcat[1][i] for subcat in others_subcategories))
            others_tuple = [("Others", tuple(others_values))]
            cleaned_data["series"] = (
                cleaned_data.get("series", [])[:max_num_categories_to_show] + others_tuple
            )

            others_subcategories_original = cleaned_data.get("original_series", [])[
                max_num_categories_to_show:
            ]
            others_values_original = []
            for i in range(len(others_subcategories_original[0][1])):
                others_values_original.append(
                    sum(subcat[1][i] for subcat in others_subcategories_original),
                )
            others_original_tuple = [("Others", tuple(others_values_original))]
            cleaned_data["original_series"] = (
                cleaned_data.get("original_series", [])[:max_num_categories_to_show]
                + others_original_tuple
            )

            for row in others_subcategories:
                subcat = row[0]
                for cat, val in zip(cleaned_data["categories"], row[1]):
                    label = cat + ", " + subcat
                    others_data.append(tuple((label, val)))

    return others_data


def create_footnote_text(cleaned_data: dict[str, Any]) -> str:
    """
    creates footnote to be shown on ppt.
     It contains the values concatenated in "Others" and values that become zero after rounding
    Args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
    Returns:
        (str): The footnote text to be shown on ppt
    """
    footnote_text = ""

    if cleaned_data.get("superscripts"):
        for superscript in cleaned_data.get("superscripts", []):
            superscript_category = superscript.get("superscript_category")

            if superscript_category == "Others":
                others_text = "; ".join(
                    [f"{cat}: {round(val):,}" for cat, val in superscript.get("original_values")],
                )
                # add newline if text exists
                if footnote_text != "":
                    footnote_text += "\n"
                footnote_text += f"{superscript.get('superscript')}: {others_text}"
            elif superscript_category == "zero_val":
                original_val = round(superscript.get("original_values")[0])
                if footnote_text != "":
                    footnote_text += "\n"
                footnote_text += f"{superscript.get('superscript')}: {original_val:,}"

    return footnote_text


def truncate_long_labels(display_labels: list[str]) -> list[str]:
    """
    Truncate the long labels to fit in the chart
    args:
        display_labels (list[str]): The labels to be shown on charts.
    returns:
        (list[str]): The truncated labels
    """
    truncated_labels = []
    for label in display_labels:
        if len(label) > ppt_generation_config["truncate_label_threshold"]:
            log.info("Label is too long. Truncating %s", label)
            label = label[: ppt_generation_config["truncate_label_threshold"]] + "..."
            truncated_labels.append(label)
        else:
            truncated_labels.append(label)
    return truncated_labels


def create_display_labels(
    cleaned_data: dict[str, Any],
    chart_type: str,
    other_values: list[Any],
) -> None:
    """
    Create the labels to be shown on charts. Rename the labels with 0 values to ">0".
     Create "superscripts" for  Others and 0 values which will be explained in footnote
    args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        chart_type (str): The type of chart to be generated based on the user's question.
    returns:
        None
    """
    display_labels = []
    superscripts = []
    superscript_cnt = 0
    if chart_type in ["bar", "pie"]:
        for val, category_label, original_val in zip(
            cleaned_data.get("series", [])[0][1],
            cleaned_data["categories"],
            cleaned_data.get("original_series", [])[0][1],
        ):
            if category_label == "Others":
                superscript_cnt += 1
                category_label = category_label + to_superscript(superscript_cnt)

                superscripts.append(
                    {
                        "category_label": category_label,
                        "original_values": other_values,
                        "superscript": superscript_cnt,
                        "superscript_category": "Others",
                    },
                )
            elif val == 0:
                superscript_cnt += 1
                category_label = category_label + to_superscript(superscript_cnt)
                original_values = [int(round(original_val, 0))]

                superscripts.append(
                    {
                        "category_label": category_label,
                        "original_values": original_values,
                        "superscript": superscript_cnt,
                        "superscript_category": "zero_val",
                    },
                )

            display_labels.append(category_label)
        if superscripts:
            cleaned_data["superscripts"] = superscripts
    else:
        display_labels = copy.deepcopy(cleaned_data["categories"])

    display_labels = truncate_long_labels(display_labels)
    cleaned_data["display_labels"] = display_labels


def prepare_data_for_visualization(
    cleaned_data: dict[str, Any],
    chart_type: str,
    rounding: str,
    chart_title: str,
) -> tuple[dict[str, Any], str, str, str, str]:
    """
    Performs rounding, soring and removes the "#" values and 0 values
    Args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        chart_type (str): The type of chart to be generated based on the user's question.
        rounding (str): The rounding key to be applied on the data.
        chart_title(str): The title of the chart to be generated.
    Returns:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        footnote (str): The footnote text to be shown on ppt
        rounding (str): The revised rounding key to be applied on the data.
        chart_title (str): The title of the chart to be generated.
        chart_type (str): Updated chart_type based on the data
    """
    llm_generated_chart_type = chart_type
    chart_type = update_chart_type_as_per_data(cleaned_data, chart_type)
    sort_filter_data(llm_generated_chart_type, cleaned_data)
    rounding, chart_title = adopt_rounding_for_chart_type(
        chart_type,
        rounding,
        cleaned_data,
        chart_title,
    )
    cleaned_data = perform_rounding(cleaned_data, rounding, chart_type)
    others_values = concatenate_to_others(cleaned_data, llm_generated_chart_type)
    create_display_labels(cleaned_data, chart_type, others_values)
    footnote = create_footnote_text(cleaned_data)
    return cleaned_data, footnote, rounding, chart_title, chart_type


def replace_dummy_pattern(data_list: list[str]) -> list[str]:
    """
    Replace the dummy pattern with the quotes
    args:
        data_list (list[str]): The list of data extracted from the response.
    returns:
        (list[str]): The list of data with dummy pattern replaced
    """
    new_list = [i.replace("&|$*", '"') for i in data_list]
    return new_list


def parse_llm_response(response: dict[str, Any]) -> tuple[dict[str, Any], str, list, str, str]:
    """
    Parse the response from the LLM model and return the cleaned data and insights.
    Args:
        response (dict[str, Any]): The response from the LLM model.
    Returns:
        tuple[dict[str, Any], str, list, str, str]:cleaned data, title, insights, title, rounding.
    """
    # Ensure cleaned data is dict
    if isinstance(response["cleaned_data"], str):
        cleaned_data = ast.literal_eval(response["cleaned_data"].strip())
    else:
        cleaned_data = response["cleaned_data"]
    cleaned_data["categories"] = replace_dummy_pattern(cleaned_data["categories"])

    # Ensure cleaned data is a list
    if isinstance(response["insights"], str):
        insights = ast.literal_eval(response["insights"].strip())
    else:
        insights = response["insights"]
    insights = replace_dummy_pattern(insights)

    title = response["title"].replace("&|$*", '"')
    chart_title = response["chart_title"].replace("&|$*", '"')
    rounding = response["rounding"]

    return cleaned_data, title, insights, chart_title, rounding


def truncate_very_long_data(data_from_answer: list[Any]) -> list[Any]:
    """
    Truncate the data if it is very long > 30 datapoints
    Args:
        data_from_answer (list[Any]): The data extracted from the response.
    Returns:
        (list[Any]): The truncated data.
    """
    try:
        data_from_answer_parsed = ast.literal_eval(data_from_answer[0][0])
    except (ValueError, IndentationError, SyntaxError):
        return data_from_answer
    max_data_limit = ppt_generation_config["truncate_dax_output_threshold"]
    if len(data_from_answer_parsed) > ppt_generation_config["truncate_dax_output_threshold"]:
        log.info(f"DAX output is too long with {len(data_from_answer_parsed)}.")
        log.info(f"Truncating to {max_data_limit}")
        data_from_answer = [tuple([str(data_from_answer_parsed[:max_data_limit])])]
    return data_from_answer


def update_chart_type_as_per_data(cleaned_data: dict[str, Any], chart_type: str) -> str:
    """
    Change line chart to bar char if data is less than covert_line_chart_to_bar_chart_threshold
    and has single category
    args:
        cleaned_data (dict[str, Any]): The processed and cleaned data extracted from the response.
        chart_type (str): chart_type generated by LLM.
    returns:
        (str): The updated chart type
    """
    # convert line chart to bar chart if data is short and has single line
    if (
        chart_type == "line"
        and len(cleaned_data["series"]) == 1
        and len(cleaned_data["categories"])
        < ppt_generation_config["covert_line_chart_to_bar_chart_threshold"]
    ):
        log.info("Converting line chart to bar chart as data is less")
        chart_type = "bar"
    return chart_type


@log_time
def generate_cleaned_data_title_insights(
    user_question: str,
    data_from_answer: list[Any],
    chart_type: str,
    currency: str,
    summarized_output: str,
) -> tuple[dict[str, Any], str, list, str, str, str, str]:
    """
    Generates cleaned data, a title, and insights based on the user's question and provided data.

    Args:
        user_question (str): The question posed by the user for which insights are to be generated.
        data_from_answer (list[Any]): Data used to generate the response.
        This could be of any type depending on the context
                          and the structure expected by the prompt generation function.
        chart_type (str): The type of chart to be generated based on the user's question.
        currency (str): The currency symbol to be used in the insights.
        summarized_output (str): The summarized output of the data.

    Returns:
        tuple[dict[str, Any], str, list, str, str, str]: A tuple containing 5 elements:
            - `cleaned_data`: The processed and cleaned data extracted from the response.
            - `title`: The title generated based on the input question and data.
            - `insights`: The insights derived from the data in response to the user's question.
            - `chart_title`: The title of the chart to be generated.
            - `rounding`: The rounding key to be applied on the data.
            - 'chart_type' : Updated chart_type based on the data

    Example:
        user_question = "What are the key trends in the sales data?"
        data_from_answer = some_data_source
        cleaned_data, title, insights = generate_cleaned_data_title_insights(
            user_question=user_question,
            data_from_answer=data_from_answer
        )
    """
    data_from_answer = truncate_very_long_data(data_from_answer)

    prompt = ppt_generate_cleaned_data_title_insights_prompt(
        data_from_answer=data_from_answer,
        user_question=user_question,
        few_shot_examples=few_shot_examples,
        chart_type=chart_type,
        currency=currency,
        summarized_output=summarized_output,
    )

    # choose the right model based on token length
    model_name = choose_model_based_on_token_length(prompt[0]["content"])

    if model_name:
        llm_repsonse = generate_chat_response(
            messages=prompt,
            model=model_name,
            response_format="json_object",
        )
        llm_repsonse = llm_repsonse.replace('\\"', "&|$*")
        try:
            response = ast.literal_eval(llm_repsonse)

        except (ValueError, IndentationError, SyntaxError):
            log.info(
                "Could not parse LLM response. Skipping PPT generation. LLM response is %s",
                str(llm_repsonse),
            )
            cleaned_data, title, insights, chart_title, rounding, footnote = (
                {"error": "llm_output_parsing_error"},
                "llm_output_parsing_error",
                ["llm_output_parsing_error"],
                "llm_output_parsing_error",
                "llm_output_parsing_error",
                "llm_output_parsing_error",
            )
            return cleaned_data, title, insights, chart_title, rounding, footnote, chart_type

        log.info(f"response from model: {response}")
        log.info(f"response from model's type is: {type(response)}")

        try:
            cleaned_data, title, insights, chart_title, rounding = parse_llm_response(response)
            (
                cleaned_data,
                footnote,
                rounding,
                chart_title,
                chart_type,
            ) = prepare_data_for_visualization(cleaned_data, chart_type, rounding, chart_title)

            log.info(
                f"cleaned_data is: {cleaned_data}, title is: {title}, "
                f"insights is: {insights}, chart title is : {chart_title}",
            )

        except (ValueError, IndentationError, SyntaxError):
            log.info("Could not parse LLM response while cleaning data. Skipping PPT generation.")

            cleaned_data, title, insights, chart_title, rounding, footnote = (
                {"error": "llm_output_parsing_error"},
                "llm_output_parsing_error",
                ["llm_output_parsing_error"],
                "llm_output_parsing_error",
                "llm_output_parsing_error",
                "llm_output_parsing_error",
            )

    else:
        log.info("Prompt exceeds token limit. Skipping PPT generation.")
        cleaned_data, title, insights, chart_title, rounding, footnote = (
            {"error": "token_limit_exceeded"},
            "token_limit_exceeded",
            ["token_limit_exceeded"],
            "token_limit_exceeded",
            "token_limit_exceeded",
            "token_limit_exceeded",
        )

    return cleaned_data, title, insights, chart_title, rounding, footnote, chart_type


def add_additional_labels(
    cleaned_data: dict[str, Any],
    slide: pptx.slide.Slide,
    chart_type: str,
) -> None:
    """
    add additional labels to the chart
    args:
        additional_labels (list[str]): The additional labels to be shown on the chart.
        slide (pptx.slide.Slide): The slide object to which the labels are to be added.
    returns:
        None
    """
    if chart_type == "stacked":
        additional_labels = []
        for i in range(len(cleaned_data["categories"])):
            additional_labels.append(sum(subcat[1][i] for subcat in cleaned_data["series"]))
        chart_start_left = Cm(2.93)
        chart_end_left = Cm(20.5)
        chart_top = Inches(3)
        bar_width = (chart_end_left - chart_start_left) / len(additional_labels)

        for i, label in enumerate(additional_labels):
            left = chart_start_left + i * bar_width
            top = chart_top
            height = Inches(0.25)
            width = Inches(0.75)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            para = text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.text = str(round(label, ppt_generation_config["round_digits_after_decimal"]))


@log_time
def generate_ppt_slide(
    chart_type: str,
    cleaned_data: dict[str, Any],
    title: str,
    insights: list,
    chart_title: str,
    template_path: str,
    user_question: str,
    category: str,
    footnote: str,
) -> bytes:
    """
    Generates a PowerPoint slide with a title, a set of bullet-point insights.

    Args:
        chart_type (str): The template for slide one of "bar", "stacked_bar", "pie" and "line"
        cleaned_data (dict[Any, Any]): The processed and cleaned data extracted from the response.
        title (str): title of the slide
        insights (str): text inside the insights box
        chart_title (str): title of the chart
        template_path (str): path of the template to use
        user_question (str): The question posed by the user for which insights are to be generated.
        category (str): The category of the material E.g. Bearing
        footnote (str): The footnote text to be shown on ppt

    Returns:
        bytes: A byte array representing the PowerPoint file containing the generated slide.

    """

    # prs = Presentation(f"{data_path}/{ppt_generation_config['template_path']}")
    prs = Presentation(template_path)
    slide = prs.slides[0]
    # Locate the title and text box elements by index
    index_title = find_index_of_element(slide, "title")
    index_key_insights = find_index_of_element(slide, "elements")
    index_key_chart_title = find_index_of_element(slide, "chart_title")
    index_key_footnote1 = find_index_of_element(slide, "footnote1")
    index_key_footnote2 = find_index_of_element(slide, "footnote2")
    index_key_category = find_index_of_element(slide, "category")

    # Add the footnote1
    footnote_box1 = slide.shapes[index_key_footnote1]
    footnote_box1.text_frame.text = footnote  # Set the footnote
    # Change the font size of the footnote to 8 points
    for paragraph in footnote_box1.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(
                8,
            )  # Set the font size to 8 points

    # Add the footnote2
    footnote_box2 = slide.shapes[index_key_footnote2]
    footnote_box2.text_frame.text = ppt_generation_config["footnote_disclaimer"]  # Set the footnote
    footnote_box2.text_frame.paragraphs[0].font.size = Pt(8)

    # Add the Category
    category_box = slide.shapes[index_key_category]
    category_box.text_frame.text = f"Category: {category}"  # Set the Category
    for paragraph in category_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)

    # Add the title
    title_box = slide.shapes[index_title]
    title_box.text_frame.text = title  # Set the title

    # Change the font size of the title to 24 points
    for paragraph in title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(
                24,
            )  # Set the font size to 24 points

    # Add the chart title
    chart_title_box = slide.shapes[index_key_chart_title]
    chart_title_box.text_frame.text = chart_title  # Set the title

    # Change the font size of the title to 18 points
    for paragraph in chart_title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(
                18,
            )  # Set the font size to 18 points
            # Add the simple sentence to the specified text box
    text_box = slide.shapes[index_key_insights]
    text_frame = text_box.text_frame

    # Clear any existing text in the text frame
    text_frame.clear()

    # Iterate over each line of insights to add them as separate bullet points
    for line in insights:
        if line.strip():  # Ensure there is text to add (ignores empty lines)
            paragraph = text_frame.add_paragraph()  # Create a new paragraph for each bullet point
            paragraph.text = line.strip()  # Set the text for the paragraph
            paragraph.level = 1  # Set bullet level (0 is the default)
            paragraph.font.size = Pt(14)  # Set font size to 12 points

    # Adjust line spacing
    for paragraph in text_frame.paragraphs:
        paragraph.line_spacing = Pt(16)  # Set line spacing to 12 points if needed

    # Modify the title
    chart_type = chart_type.lower()

    populate_charts(slide, cleaned_data, chart_type)

    add_additional_labels(cleaned_data, slide, chart_type)
    # Save the presentation to a byte buffer
    buffer = BytesIO()
    prs.save(buffer)

    log.info("user question: %s", user_question)
    # write to file
    if os.path.exists("downloads"):
        prs.save("downloads/" + str(user_question) + ".pptx")

    buffer.seek(0)  # Rewind the buffer
    file_bytes = buffer.getvalue()  # write bytes array in dev

    # write bytes array in dev
    return file_bytes


@log_time
def run_ppt_generation(
    json_file: str,
) -> dict[str, Any]:
    """
    Run the complete process of generating a
    PowerPoint presentation based on the provided JSON data.

    Args:
    json_file (str): Input JSON data as a string.

    Returns:
    dict[str, Any]: The generated PowerPoint file along with other metadata or an error response.
    """
    # Step 1: Query the data and user's question
    user_question, data_from_answer, currency, summarized_output, category, tenant_id = get_data(
        json_file,
    )
    set_tenant_and_intent_flow_context(tenant_id, UseCase.PPT_GENERATION)

    log.info("Query Result: %s", data_from_answer)

    # Step 2: Get the chart type
    chart_type = generate_chart_type(user_question, data_from_answer)
    log.info("chart type: %s", chart_type)

    # Step 3: Generate the cleaned data, title, and insights

    if chart_type is None:
        # The data is too much and does not fit in the token limit of the LLM
        log.info("Too much data while inferring chart type. Skipping PPT generation.")
        response_data = {
            "file_name": "failed_presentation.pptx",
            "file_type": "pptx",
            # "file_data": base64.b64encode(ppt_file).decode("utf8"),
            "encoding_type": "base64",
            "char_encoding": "utf8",
            "is_encoded": True,
            # "chat_id": chat_id,
            "message": "PPT not generated due to too much data",
            "response_type": "data_issue",
        }

    elif chart_type == "Invalid":
        # data is too less to generate PPT
        log.info("Too less data. Skipping PPT generation.")
        response_data = {
            "file_name": "failed_presentation.pptx",
            "file_type": "pptx",
            "encoding_type": "base64",
            "char_encoding": "utf8",
            "is_encoded": True,
            "message": "PPT not generated due to insufficient data",
            "response_type": "data_issue",
        }
    else:
        (
            cleaned_data,
            title,
            insights,
            chart_title,
            rounding,
            footnote,
            chart_type,
        ) = generate_cleaned_data_title_insights(
            data_from_answer=data_from_answer,
            user_question=user_question,
            chart_type=chart_type,
            currency=currency,
            summarized_output=summarized_output,
        )
        log.info(f"cleaned_data: {cleaned_data}")
        log.info(f"title: {cleaned_data}")
        log.info(f"insights: {insights}")
        log.info(f"chart_title: {chart_title}")
        log.info(f"rounding: {rounding}")
        log.info(f"footnote: {footnote}")

        if cleaned_data == {"error": "token_limit_exceeded"}:
            # The data is too much
            response_data = {
                "file_name": "failed_presentation.pptx",
                "file_type": "pptx",
                "encoding_type": "base64",
                "char_encoding": "utf8",
                "is_encoded": True,
                "message": "PPT not generated due to too much data",
                "response_type": "data_issue",
            }
        elif cleaned_data == {"error": "llm_output_parsing_error"}:
            response_data = {
                "file_name": "failed_presentation.pptx",
                "file_type": "pptx",
                # "file_data": base64.b64encode(ppt_file).decode("utf8"),
                "encoding_type": "base64",
                "char_encoding": "utf8",
                "is_encoded": True,
                # "chat_id": chat_id,
                "message": "PPT not generated due to an internal error. Please try again.",
                "response_type": "data_issue",
            }
        else:
            # Continue with further processing using cleaned_data, title, and insights as needed
            # For example:
            log.info("Cleaned Data:  %s", cleaned_data)
            log.info("Title:  %s", title)
            log.info("Insights:  %s", insights)
            log.info("Chart title:  %s", chart_title)
            log.info("Rounding:  %s", rounding)

            # Step 4: Import the right template
            template_path = get_pptx_template_path(chart_type)
            log.info("Importing file: %s", template_path)

            # Step 5: Populate the slide
            ppt_file = generate_ppt_slide(
                template_path=template_path,
                chart_type=chart_type,
                cleaned_data=cleaned_data,
                insights=insights,
                title=title,
                chart_title=chart_title,
                user_question=user_question,
                category=category,
                footnote=footnote,
            )

            log.info("Generating PPT ended")

            if not ppt_file:
                raise ValueError("PPT file was not generated.")
            print(f"Returning PPT file, length: {len(ppt_file)}")

            # Return the PowerPoint file as bytes
            response_data = {
                "file_name": "presentation.pptx",
                "file_type": "pptx",
                "file_data": base64.b64encode(ppt_file).decode("utf8"),
                "encoding_type": "base64",
                "char_encoding": "utf8",
                "is_encoded": True,
                "message": "PPT generated",
                "response_type": "general",
            }
    return response_data
