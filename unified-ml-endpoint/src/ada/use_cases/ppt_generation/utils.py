"""PPT generation use case utils"""

import json
import os
from pathlib import Path
from typing import Any

import pptx.slide
import tiktoken
from pptx.chart.data import CategoryChartData

from ada.components.db.pg_connector import PGConnector
from ada.use_cases.key_facts.configuration import Configuration
from ada.utils.config.config_loader import read_config
from ada.utils.logs.logger import get_logger

log = get_logger("ppt-generation-v1")
ppt_generation_config = read_config("use-cases.yml")["ppt_generation"]
model_conf = read_config("models.yml")


def find_shape_index_by_name(slide: pptx.slide.Slide, shape_name: str) -> int:
    """
    Finds the index of the shape with the given name on the provided slide.

    Parameters:
        slide (pptx.slide.Slide): A slide object from the PowerPoint presentation.
        shape_name (str): The name of the shape to find (case-insensitive).

    Returns:
        (int): The index of the shape if found, otherwise -1.
    """
    for i, shape in enumerate(slide.shapes):
        if shape_name.lower() in str(shape.shape_type).lower():
            return i
    return -1  # Return -1 if the shape with the given name is not found


def find_index_of_element(slide: pptx.slide.Slide, element_text: str) -> int | None:
    """
    Find the index of an element in a slide
    Args:
        slide (pptx.slide.Slide): A slide object from the PowerPoint presentation.
        element_text (str): The text of the element to find.
    Returns:
        (int | None): element index that matches element_text
    """
    elements = [
        i
        for i, shape in enumerate(slide.shapes)
        if shape.has_text_frame and shape.text_frame.text == element_text
    ]
    return_index = None
    if len(elements) == 1:
        return_index = elements[0]
    return return_index


def get_data(json_file: str) -> tuple[str, list[Any], str, str, str, str]:
    """
    Get the user question and answer data from the database based on the provided JSON data.

    Args:
    json_file (str): Input JSON data as a string.

    Returns:
    (tuple[str, list[Any], str, str, str]): A tuple containing the user_question, data_from_answer,
        currency, summarized_output, category
    """

    json_data = json.loads(json_file)
    tenant_id = json_data.get("tenant_id", "")
    category = json_data.get("category", "")
    request_id = json_data.get("request_id", "")
    query_request_id = json_data.get("query_request_id", "")
    tenant_id = json_data.get("tenant_id", "")
    category = json_data.get("category", "")
    request_id = json_data.get("request_id", "")
    query_request_id = json_data.get("query_request_id", "")

    log.info("Request ID: %s\nUser query: %s\n", request_id, query_request_id)

    # Connect to DB to retrieve the user's question and answer
    pg_db_conn = PGConnector(tenant_id=tenant_id)
    # dev db (deploy the endpoint , local db

    user_question = pg_db_conn.select_records_with_filter(
        table_name="dax_queries",
        filtered_columns=["user_question"],
        filter_condition=(
            f" request_id = '{query_request_id}' "
            f"and user_category = '{category}' "
            "and execution_status=200"
        ),
    )
    user_question = user_question[0][0]
    # table name, and column name
    log.info("Query: %s", user_question)

    data_from_answer = pg_db_conn.select_records_with_filter(
        table_name="dax_queries",  # unified_chat_history
        filtered_columns=["execution_output"],
        filter_condition=(
            f" request_id = '{query_request_id}' "
            f"and user_category = '{category}' "
            "and execution_status=200"
        ),
    )

    summarized_output = pg_db_conn.select_records_with_filter(
        table_name="dax_queries",  # unified_chat_history
        filtered_columns=["summarised_output"],
        filter_condition=(
            f" request_id = '{query_request_id}' "
            f"and user_category = '{category}' "
            "and execution_status=200"
        ),
    )
    summarized_output = summarized_output[0][0]

    try:
        category_configuration = Configuration(required_configuration_keys={"currency"})
        category_configuration.load_configuration(pg_connector=pg_db_conn, category=category)
        currency = (
            category_configuration.configuration_dict.get("currency", {}).get("currency", "EUR")
            or "EUR"
        )
    except ValueError:
        log.info(
            "Current category does not have currency configuration. Using default currency: EUR",
        )
        currency = "EUR"

    return user_question, data_from_answer, currency, summarized_output, category, tenant_id


def get_pptx_template_path(chart_type: str | None) -> str | None:
    """
    Get the path to the PowerPoint template file based on the chart type.

    Args:
        chart_type (str): Chart Type selected in previous step.

    Returns:
        str or None: The path to the matching PowerPoint file or None if no match is found.

    """
    templates_folder_path = Path(Path(__file__).parents[4], "data/ppt_generation/templates")
    log.info(f"templates folder path: {templates_folder_path}")
    template_path = os.path.join(templates_folder_path, f"template_{chart_type}.pptx")
    return template_path


def populate_charts(slide: pptx.slide.Slide, cleaned_data: dict[str, Any], chart_type: str) -> None:
    """
    Populates a stacked chart on the given slide with the provided data.

    This function finds the chart on the slide by its predefined name ("CHART"),
    and then replaces its data with the cleaned data provided. The chart is updated
    with new categories and series values, and its legend is configured accordingly.

    Args:
        slide (pptx.slide.Slide): The slide containing the chart to be populated.
        cleaned_data (dict[ str, Any]): A dictionary containing the data to be used for the chart.
        cleaned_data (dict[ str, Any]): A dictionary containing the data to be used for the chart.
            It should contain:
                - "categories" (list): A list of categories for the chart.
                - "series" (list of tuples): Each tuple contains a series name and a list
                  of corresponding values.
        chart_type (str): The type of chart to be populated.
                - Can be one of : "line", "bar", "pie", "stacked"
        chart_type (str): The type of chart to be populated.
                - Can be one of : "line", "bar", "pie", "stacked"

    Example:
        cleaned_data = {
            "categories": ["Category 1", "Category 2", "Category 3"],
            "series": [
                ("Series 1", [10, 20, 30]),
                ("Series 2", [15, 25, 35]),
            ],
        }
        populate_charts(slide, cleaned_data, "stacked")
        populate_charts(slide, cleaned_data, "stacked")

    Returns:
        None
    """
    chart_index = find_shape_index_by_name(slide, "CHART")

    chart = slide.shapes[chart_index].chart  # Accessing the chart directly by the correct index
    chart_data = CategoryChartData()
    chart_data.categories = cleaned_data["display_labels"]
    if chart_type != "bar":
        for series_name, series_values in cleaned_data["series"]:
            chart_data.add_series(series_name, series_values)
    else:
        for series_name, series_values in cleaned_data["series"]:
            if isinstance(series_values, tuple):
                log.info(f"series is %s: {series_values}")
                chart_data.add_series(series_name, series_values)
            else:
                series_values_tuple = tuple([series_values])
                log.info(f"series is %s: {series_values}")
                chart_data.add_series(series_name, series_values_tuple)

    chart.replace_data(chart_data)
    if chart_type in ["stacked", "pie"]:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.series[0].smooth = True
    elif chart_type == "bar":
        chart.has_legend = False
        chart.series[0].smooth = True
    elif chart_type == "line":
        chart.has_legend = True
        chart.legend.include_in_layout = False


def get_max_token_length_of_model(model_name: str) -> int:
    """
    Get the max token length of the input model name based on the model configuration file.

    args:
    model_name (str): The name of the model to get the max token length for.
    model_name (str): The name of the model to get the max token length for.

    returns:
    int: The max token length of the model. If the model name is not found, return None.
    """
    max_token_length = 0
    for model in model_conf["models"]:
        if model["model_name"] == model_name:
            max_token_length = model["max_tokens"]
    return max_token_length


def choose_model_based_on_token_length(prompt: str) -> str | None:
    """
    Choose the right model based on the token length of the user question and data.
    If the token length is greater than 32k then return None.
    Args:
        prompt (str): The prompt to be used to generate the response.
    Args:
        prompt (str): The prompt to be used to generate the response.

    Returns:
        (str): The name of the model to be used to generate the response.
    """
    # use tiktoken to get the token length
    encoding = tiktoken.get_encoding("cl100k_base")
    num_tokens = len(encoding.encode(prompt))
    log.info("Prompt length in tokens: %s", num_tokens)

    primary_model_token_limit = get_max_token_length_of_model(
        ppt_generation_config["model_name_primary"],
    )
    secondary_model_token_limit = get_max_token_length_of_model(
        ppt_generation_config["model_name_secondary"],
    )

    if num_tokens <= int(primary_model_token_limit):
        selected_model_name = ppt_generation_config["model_name_primary"]
    elif num_tokens <= (secondary_model_token_limit or -1):
        selected_model_name = ppt_generation_config["model_name_secondary"]
    else:
        log.info("Too long prompt of token length: %s", num_tokens)
        selected_model_name = None
    return selected_model_name


def to_superscript(number: int) -> str:
    """
    Convert a number to a superscript string.
    args:
    number (int): The number to convert to superscript.
    returns:
    str: The number converted to a superscript string.
    """
    superscript_map = {
        "0": "⁰",
        "1": "¹",
        "2": "²",
        "3": "³",
        "4": "⁴",
        "5": "⁵",
        "6": "⁶",
        "7": "⁷",
        "8": "⁸",
        "9": "⁹",
    }

    return "".join(superscript_map[digit] for digit in str(number))
